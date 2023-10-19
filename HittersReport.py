from ast import walk
from cgi import print_directory
from email.mime import image
import math
from cmath import nan
from unicodedata import name
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from tkinter import filedialog
import seaborn as sns
import os
from pptx import Presentation
from pptx.util import Pt
import comtypes.client
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from scipy.interpolate import interp2d
from scipy.interpolate import RectBivariateSpline
from scipy.ndimage.filters import gaussian_filter
from scipy import interpolate
import matplotlib.patches as patches
from PIL import Image
from pptx.dml.color import RGBColor
import matplotlib.lines as lines


##Create a report that takes in hitter data from a CSV file with muliple
##Trackman games and generated a PPTX and PDF containing charts and other 
##Data from the hitter

#Asks for CSV File
csv_file = filedialog.askopenfilename()
count = 0
#csv_file = "C:\\Users\\cmhea\\OneDrive\\Documents\\baseball\\Knights\\KnightsCSVs2023\\Verified\\Verfied620.csv"
csv_df = pd.read_csv(csv_file)
csv_df = csv_df[csv_df['BatterTeam'] == 'COR_KNI']
names = ['Quinn, Tyler','Hedges, Ethan','Yukumoto, Ty','Hott, Ethan','Stem, Samuel','Shimao, Tate','Avila, Blake','Jones, Merit','Aroz, Anson','Le, Mason','Howard, Tyler','Call, Phoenix','Ng, JC']

#Methods Job is to Clean Up the csv data frame to one only containg rows 
#of player we want and collumns we need for the swing density chart
def get_player_df():
    player_df = csv_df.drop(csv_df[csv_df.Batter != names[count]].index)
    player_df = player_df.reset_index(drop=True)
    return player_df

def all_results(player_df):
    df1 = player_df[(player_df['PitchCall'].isin(['InPlay', 'HitByPitch'])) | (player_df['KorBB'].isin(['Strikeout', 'Walk']))]
    df1.to_csv(os.path.join("Sheets", names[count], names[count] + '.csv'))
    

def csv_to_swing_df():
    global csv_df
    player_df = csv_df

    #Drops all rows where player isn't hitting
    player_df = player_df.drop(player_df[player_df.Batter != names[count]].index)
    #Drops all rows where player didn't swing
    player_df = player_df.drop(player_df[player_df.PitchCall == 'BallCalled'].index)
    player_df = player_df.drop(player_df[player_df.PitchCall == 'StrikeCalled'].index)
    player_df = player_df.drop(player_df[player_df.PitchCall == 'BallinDirt'].index)
    player_df = player_df.drop(player_df[player_df.PitchCall == 'HitByPitch'].index)
    #Creating a list of collums to remove from DF and then removing collums we need from list
    remove_list = player_df.columns.values.tolist()
    remove_list.remove("Batter")
    remove_list.remove("PlateLocHeight")
    remove_list.remove("PlateLocSide")
    remove_list.remove("PitchCall")
    remove_list.remove("BatterSide")
    #Removes all collums other than those with .remove above from data frame
    player_df = player_df.drop(remove_list, axis=1)
    #Switches plate loc side values to be in catcher view by multiplying by -1
    player_df['PlateLocSide'] = (player_df['PlateLocSide'] * -1)
    print(player_df)


    return player_df

# Currently Unused but can replace the swing df to make a whiff location chart
def csv_to_whiff_df():
    global csv_df
    player_df = csv_df

    #Drops all rows where player isn't hitting
    player_df = player_df.drop(player_df[player_df.Batter != names[count]].index)
    #Drops all rows where player didn't whiff
    player_df = player_df.drop(player_df[player_df.PitchCall == 'BallCalled'].index)
    player_df = player_df.drop(player_df[player_df.PitchCall == 'StrikeCalled'].index)
    player_df = player_df.drop(player_df[player_df.PitchCall == 'BallinDirt'].index)
    player_df = player_df.drop(player_df[player_df.PitchCall == 'FoulBall'].index)
    player_df = player_df.drop(player_df[player_df.PitchCall == 'InPlay'].index)
    #Creating a list of collums to remove from DF and then removing collums we need from list
    remove_list = player_df.columns.values.tolist()
    remove_list.remove("Batter")
    remove_list.remove("PlateLocHeight")
    remove_list.remove("PlateLocSide")
    remove_list.remove("PitchCall")
    remove_list.remove("BatterSide")
    #Removes all collums other than those with .remove above from data frame
    player_df = player_df.drop(remove_list, axis=1)
    #Switches plate loc side values to be in catcher view by multiplying by -1
    player_df['PlateLocSide'] = (player_df['PlateLocSide'] * -1)
    print(player_df)


    return player_df

#Takes in a dataframe and creates a seaborn 2d density plot based on pitch location of dataframe
def swing2d_density_plot(player_df):
    #Pulls image for background will have to imput if statemnt to deptermine right vs left
    rhhs = ['Hedges, Ethan','Hott, Ethan','Howard, Tyler','Krieg, Jacob','Quinn, Tyler','Rowe, Cameron','Schoppe, Stanley','Segel, Kellen',
         'Shimao, Tate','Stone, Jonathan','Thiele, Luke', 'Le, Mason','Call, Phoenix','Shimao, Tate']
    if names[count] in rhhs:
        img = plt.imread("RHH.png")
    else:
        img = plt.imread('LHH.png')
    fig, ax = plt.subplots(figsize=(6, 6))
    sns.set_style("white")
    #Creates density plot, camp is color scheme and alpha is transperacy
    chart = sns.kdeplot(x=player_df.PlateLocSide, y=player_df.PlateLocHeight,cmap='rocket_r', shade=True, bw_adjust=.55, ax=ax, alpha = 0.65)
    #creates demenstions for graph plus displays image
    ax.imshow(img, extent=[-2.63,2.665,-0.35,5.30], aspect=1)

    chart.set(xticklabels=[])  
    chart.set(xlabel=None)
    chart.tick_params(bottom=False)
    chart.set(yticklabels=[])  
    chart.set(ylabel=None)
    chart.tick_params(left=False)   # remove the ticks
    #rect = patches.Rectangle((-0.708333, 1.6466667), 1.4166667, 1.90416667, linewidth=1, edgecolor='black', facecolor='none')
    #ax.add_patch(rect)
    #creates path for plot to be saved in there isn't one
    newpath = os.path.join("Sheets", names[count])
    if not os.path.exists(newpath):
        os.makedirs(newpath)
    #saves plot in folder
    plt.axis('off')
    plt.savefig(os.path.join("Sheets", names[count], 'swingchart.png'),bbox_inches='tight', pad_inches = 0)


    return

#Gets Metrics for both tables
def find_table_metrics():
    global csv_df
    player_df = csv_df
    #Drops all rows where player isn't hitting
    player_df = player_df.drop(player_df[player_df.Batter != names[count]].index)
    remove_list = player_df.columns.values.tolist()
    remove_list.remove("Batter")
    remove_list.remove("PlateLocHeight")
    remove_list.remove("PlateLocSide")
    remove_list.remove("PitchCall")
    remove_list.remove("BatterSide")
    remove_list.remove("KorBB")
    remove_list.remove("PlayResult")
    remove_list.remove("ExitSpeed")
    remove_list.remove("TaggedHitType")
    remove_list.remove("RunsScored")

    #Removes all collums other than those with .remove above from data frame
    player_df = player_df.drop(remove_list, axis=1)
    avg_ev_df = player_df.drop(player_df[player_df.ExitSpeed < 55].index)
    avg_ev = "%.1f" % round(avg_ev_df["ExitSpeed"].mean(),1)
    max_ev = "%.1f" % round(avg_ev_df["ExitSpeed"].max(),1)
    
    #Swing %
    swings = (player_df["PitchCall"] == "InPlay").sum() + (player_df["PitchCall"] == "FoulBall").sum() + (player_df["PitchCall"] == "StrikeSwinging").sum()
    takes = (player_df["PitchCall"] == "BallCalled").sum() + (player_df["PitchCall"] == "StrikeCalled").sum()
    swing_rate = "%.1f" % round(100*(swings/(swings+takes)),1)

    #Chase Rate
    out_of_zone_df = player_df
    indexzone  = out_of_zone_df[ (out_of_zone_df['PlateLocSide'] < 0.83083) & (out_of_zone_df['PlateLocSide'] > -0.83083) & (out_of_zone_df['PlateLocHeight'] < 3.67333) & (out_of_zone_df['PlateLocHeight'] > 1.52417)].index
    out_of_zone_df = out_of_zone_df.drop(indexzone)
    chases = (out_of_zone_df["PitchCall"] == "InPlay").sum() + (out_of_zone_df["PitchCall"] == "FoulBall").sum() + (out_of_zone_df["PitchCall"] == "StrikeSwinging").sum()
    takes_out_of_zone =  (out_of_zone_df["PitchCall"] == "BallCalled").sum() + (out_of_zone_df["PitchCall"] == "StrikeCalled").sum() + (out_of_zone_df["PitchCall"] == "HitByPitch").sum() + (out_of_zone_df["PitchCall"] == "BallinDirt").sum()
    chase_rate = "%.1f" % round(100*(chases/(chases+takes_out_of_zone)),1)

    #InZoneSwing%
    in_zone_df = player_df[(player_df['PlateLocSide'] < 0.83083) & (player_df['PlateLocSide'] > -0.83083) & (player_df['PlateLocHeight'] < 3.67333) & (player_df['PlateLocHeight'] > 1.52417)]
    in_zone_swings = (in_zone_df["PitchCall"] == "InPlay").sum() + (in_zone_df["PitchCall"] == "FoulBall").sum() + (in_zone_df["PitchCall"] == "StrikeSwinging").sum()
    in_zone_takes = (in_zone_df["PitchCall"] == "BallCalled").sum() + (in_zone_df["PitchCall"] == "StrikeCalled").sum() + (in_zone_df["PitchCall"] == "HitByPitch").sum() + (in_zone_df["PitchCall"] == "BallinDirt").sum()
    in_zone_swing_percentage = "%.1f" % round(100*(in_zone_swings/(in_zone_swings+in_zone_takes)),1)


    #K Rate
    walk_df = player_df.drop(player_df[player_df.PitchCall == 'HitByPitch'].index)
    strikeouts = (player_df["KorBB"] == "Strikeout").sum()
    plate_apearences = (walk_df["KorBB"] == "Walk").sum() + (player_df["KorBB"] == "Strikeout").sum() + (player_df["PitchCall"] == "InPlay").sum() + (player_df["PitchCall"] == "HitByPitch").sum()
    k_rate = "%.1f" % round(100*(strikeouts/plate_apearences),1)

    #(BB+HBP)/K
    walks = (walk_df["KorBB"] == "Walk").sum()
    hbps = (player_df["PitchCall"] == "HitByPitch").sum()
    bb_hbp_over_ks = "%.2f" % round(((walks + hbps)/strikeouts),2)

    #BABIP
    bip = (player_df["PitchCall"] == "InPlay").sum() - (player_df["PlayResult"] == "HomeRun").sum()
    hits_no_hrs = (player_df["PlayResult"] == "Single").sum() + (player_df["PlayResult"] == "Double").sum() + (player_df["PlayResult"] == "Triple").sum()
    babip = "%.3f" % round(hits_no_hrs/bip,3)

    #AVG
    ABs = (player_df["PitchCall"] == "InPlay").sum() - (player_df["PlayResult"] == "Sacrifice").sum() + (player_df["KorBB"] == "Strikeout").sum()
    hits = (player_df["PlayResult"] == "Single").sum() + (player_df["PlayResult"] == "Double").sum() + (player_df["PlayResult"] == "Triple").sum() + (player_df["PlayResult"] == "HomeRun").sum()
    avg = "%.3f" % round(hits/ABs,3)

    #RBIs
    inplay_df = player_df.drop(player_df[player_df.PitchCall != 'InPlay'].index)
    walk_df = player_df.drop(player_df[player_df.KorBB != 'Walk'].index)
    hbp_df = player_df.drop(player_df[player_df.PitchCall != 'HitByPitch'].index)
    rbidf = pd.concat([inplay_df,walk_df,hbp_df])
    rbi = "%.0f" % round(rbidf["RunsScored"].sum(),0)

    #Doubles, Triples, Home Runs
    doubles = (player_df["PlayResult"] == "Double").sum()
    triples = (player_df["PlayResult"] == "Triple").sum()
    homers = (player_df["PlayResult"] == "HomeRun").sum()
    
    #OBP
    obp = "%.3f" % round((walks + hbps + hits)/plate_apearences,3)

    #SLG
    slg = "%.3f" % round((hits + doubles + triples*2 + homers*3)/ABs,3)

    #OPS
    ops = "%.3f" % (round((walks + hbps + hits)/plate_apearences,3) + round((hits + doubles + triples*2 + homers*3)/ABs,3))

    #WOBA
    woba = "%.3f" % (round((0.689*walks + 0.720*hbps + 0.844*(hits-doubles-triples-homers) + 1.261*doubles + 1.601*triples + 2.072*homers)/plate_apearences,3))
    league_woba = 0.351


    #Compiles Data into one list to pass to the presentation fucntion
    data_to_pass_to_presentation = []
    data_to_pass_to_presentation.append(str(avg_ev))
    data_to_pass_to_presentation.append(str(max_ev))
    data_to_pass_to_presentation.append(str(swing_rate) + "%")
    data_to_pass_to_presentation.append(str(chase_rate) + "%")
    data_to_pass_to_presentation.append(str(in_zone_swing_percentage) + "%")
    data_to_pass_to_presentation.append(str(k_rate) + "%")
    data_to_pass_to_presentation.append(str(bb_hbp_over_ks))
    data_to_pass_to_presentation.append(str(babip).lstrip('0'))
    data_to_pass_to_presentation.append(str(woba).lstrip('0'))
    data_to_pass_to_presentation.append(str(plate_apearences))
    data_to_pass_to_presentation.append(str(ABs))
    data_to_pass_to_presentation.append(str(hits))
    data_to_pass_to_presentation.append(str(doubles))
    data_to_pass_to_presentation.append(str(triples))
    data_to_pass_to_presentation.append(str(homers))
    data_to_pass_to_presentation.append(str(rbi))
    data_to_pass_to_presentation.append(str(walks))
    data_to_pass_to_presentation.append(str(hbps))
    data_to_pass_to_presentation.append(str(strikeouts))
    data_to_pass_to_presentation.append(str(avg).lstrip('0'))
    data_to_pass_to_presentation.append(str(obp).lstrip('0'))
    data_to_pass_to_presentation.append(str(slg).lstrip('0'))
    data_to_pass_to_presentation.append(str(ops).lstrip('0'))


    print(data_to_pass_to_presentation)
    return data_to_pass_to_presentation

def find_table_metrics_using_online(stats):
    global csv_df
    stats.to_csv('stats.csv')
    player_df = csv_df
    #Drops all rows where player isn't hitting
    player_df = player_df.drop(player_df[player_df.Batter != names[count]].index)
    remove_list = player_df.columns.values.tolist()
    remove_list.remove("Batter")
    remove_list.remove("PlateLocHeight")
    remove_list.remove("PlateLocSide")
    remove_list.remove("PitchCall")
    remove_list.remove("BatterSide")
    remove_list.remove("KorBB")
    remove_list.remove("PlayResult")
    remove_list.remove("ExitSpeed")
    remove_list.remove("TaggedHitType")
    remove_list.remove("RunsScored")

    #Removes all collums other than those with .remove above from data frame
    player_df = player_df.drop(remove_list, axis=1)
    avg_ev_df = player_df.drop(player_df[player_df.ExitSpeed < 55].index)
    avg_ev = "%.1f" % round(avg_ev_df["ExitSpeed"].mean(),1)
    max_ev = "%.1f" % round(avg_ev_df["ExitSpeed"].max(),1)
    
    #Swing %
    swings = (player_df["PitchCall"] == "InPlay").sum() + (player_df["PitchCall"] == "FoulBall").sum() + (player_df["PitchCall"] == "StrikeSwinging").sum()
    takes = (player_df["PitchCall"] == "BallCalled").sum() + (player_df["PitchCall"] == "StrikeCalled").sum()
    swing_rate = "%.1f" % round(100*(swings/(swings+takes)),1)

    #Chase Rate
    out_of_zone_df = player_df
    indexzone  = out_of_zone_df[ (out_of_zone_df['PlateLocSide'] < 0.83083) & (out_of_zone_df['PlateLocSide'] > -0.83083) & (out_of_zone_df['PlateLocHeight'] < 3.67333) & (out_of_zone_df['PlateLocHeight'] > 1.52417)].index
    out_of_zone_df = out_of_zone_df.drop(indexzone)
    out_of_zone_df.dropna(subset=['PlateLocHeight', 'PlateLocSide'], inplace=True)
    chases = (out_of_zone_df["PitchCall"] == "InPlay").sum() + (out_of_zone_df["PitchCall"] == "FoulBall").sum() + (out_of_zone_df["PitchCall"] == "StrikeSwinging").sum()
    takes_out_of_zone =  (out_of_zone_df["PitchCall"] == "BallCalled").sum() + (out_of_zone_df["PitchCall"] == "StrikeCalled").sum() + (out_of_zone_df["PitchCall"] == "HitByPitch").sum() + (out_of_zone_df["PitchCall"] == "BallinDirt").sum()
    chase_rate = "%.1f" % round(100*(chases/(chases+takes_out_of_zone)),1)

    #InZoneSwing%
    in_zone_df = player_df[(player_df['PlateLocSide'] < 0.83083) & (player_df['PlateLocSide'] > -0.83083) & (player_df['PlateLocHeight'] < 3.67333) & (player_df['PlateLocHeight'] > 1.52417)]
    in_zone_df.dropna(subset=['PlateLocHeight', 'PlateLocSide'], inplace=True)
    in_zone_swings = (in_zone_df["PitchCall"] == "InPlay").sum() + (in_zone_df["PitchCall"] == "FoulBall").sum() + (in_zone_df["PitchCall"] == "StrikeSwinging").sum()
    in_zone_takes = (in_zone_df["PitchCall"] == "BallCalled").sum() + (in_zone_df["PitchCall"] == "StrikeCalled").sum() + (in_zone_df["PitchCall"] == "HitByPitch").sum() + (in_zone_df["PitchCall"] == "BallinDirt").sum()
    in_zone_swing_percentage = "%.1f" % round(100*(in_zone_swings/(in_zone_swings+in_zone_takes)),1)

    walks = int(stats.loc[0,'BB'])
    strikeouts = int(stats.loc[0,'SO'])
    hbps = int(stats.loc[0,'HBP'])
    ABs = stats.loc[0,'AB']
    bb_hbp_over_ks = "%.2f" % round((walks+hbps)/strikeouts,2)
    doubles = int(stats.loc[0,'2B'])
    triples = int(stats.loc[0,'3B'])
    homers = int(stats.loc[0,'HR'])
    hits = int(stats.loc[0,'H'])
    babip = "%.3f" % round((hits - homers)/(ABs - homers - strikeouts + stats.loc[0,'SF']),3)
    plate_apearences = ABs + stats.loc[0,'SF'] + stats.loc[0,'SH'] + walks + hbps
    
    k_rate = "%.1f" % round(float(100*(strikeouts/plate_apearences)),1)

    rbi = int(stats.loc[0,'RBI'])
    runs = int(stats.loc[0,'R'])

    avg = "%.3f" % round(hits/ABs,3)

    slg = "%.3f" % round(((hits + doubles + triples*2 + homers*3)/ABs),3)

    obp = "%.3f" % round((hits+walks+hbps)/plate_apearences,3)

    ops = "%.3f" % round(((hits + doubles + triples*2 + homers*3)/ABs)+((hits+walks+hbps)/plate_apearences),3)

    '''
    #K Rate
    k_rate = "%.1f" % round(float(stats.loc[0,'so%']),1)


    #(BB+HBP)/K
    walks = int(stats.loc[0,'bb'])
    hbps = int(stats.loc[0,'hbp'])
    strikeouts = int(stats.loc[0,'so'])
    bb_hbp_over_ks = "%.2f" % round((walks+hbps)/strikeouts,2)


    #BABIP
    babip = float(stats.loc[0,'babip'])


    #AVG
    avg = stats.loc[0,'avg']


    #RBIs
    rbi = int(stats.loc[0,'rbi'])
    runs = int(stats.loc[0,'r'])


    #Doubles, Triples, Home Runs
    doubles = int(stats.loc[0,'2b'])
    triples = int(stats.loc[0,'3b'])
    homers = int(stats.loc[0,'hr'])
    hits = int(stats.loc[0,'h'])

    #PAs & ABs
    plate_apearences = int(stats.loc[0,'pa'])
    ABs = int(stats.loc[0,'ab'])

    
    #OBP
    obp = stats.loc[0,'obp']


    #SLG
    slg = stats.loc[0,'slg']


    #OPS
    ops = stats.loc[0,'ops']
    '''


    #WOBA
    woba = "%.3f" % (round((0.689*walks + 0.720*hbps + 0.844*(hits-doubles-triples-homers) + 1.261*doubles + 1.601*triples + 2.072*homers)/plate_apearences,3))



    #Compiles Data into one list to pass to the presentation fucntion
    data_to_pass_to_presentation = []
    data_to_pass_to_presentation.append(str(avg_ev))
    data_to_pass_to_presentation.append(str(max_ev))
    data_to_pass_to_presentation.append(str(swing_rate) + "%")
    data_to_pass_to_presentation.append(str(chase_rate) + "%")
    data_to_pass_to_presentation.append(str(in_zone_swing_percentage) + "%")
    data_to_pass_to_presentation.append(str(k_rate) + "%")
    data_to_pass_to_presentation.append(str(bb_hbp_over_ks))
    data_to_pass_to_presentation.append(str(babip).lstrip('0'))
    data_to_pass_to_presentation.append(str(woba).lstrip('0'))
    data_to_pass_to_presentation.append(str(plate_apearences))
    data_to_pass_to_presentation.append(str(ABs))
    data_to_pass_to_presentation.append(str(hits))
    data_to_pass_to_presentation.append(str(doubles))
    data_to_pass_to_presentation.append(str(triples))
    data_to_pass_to_presentation.append(str(homers))
    data_to_pass_to_presentation.append(str(rbi))
    data_to_pass_to_presentation.append(str(runs))
    data_to_pass_to_presentation.append(str(walks))
    data_to_pass_to_presentation.append(str(hbps))
    data_to_pass_to_presentation.append(str(strikeouts))
    data_to_pass_to_presentation.append(str(avg).lstrip('0'))
    data_to_pass_to_presentation.append(str(obp).lstrip('0'))
    data_to_pass_to_presentation.append(str(slg).lstrip('0'))
    data_to_pass_to_presentation.append(str(ops).lstrip('0'))


    print(data_to_pass_to_presentation)
    return data_to_pass_to_presentation

#Creates data Frame for EV Heat map Catcher View
def data_frame_for_damage_chart():
    global csv_df
    damage_df = csv_df

    #Drops rows without listed hitter
    damage_df = damage_df.drop(damage_df[damage_df.Batter != names[count]].index)

    #Remove unnesassary columns
    remove_list = damage_df.columns.values.tolist()
    remove_list.remove("Batter")
    remove_list.remove("PlateLocHeight")
    remove_list.remove("PlateLocSide")
    remove_list.remove("BatterSide")
    remove_list.remove("ExitSpeed")
    damage_df = damage_df.drop(remove_list, axis=1)

    #Removes Rows where EV is nan
    drop_index = []
    for index, row in damage_df.iterrows():
        if math.isnan(row['ExitSpeed']) == True:
            drop_index.append(index)
    damage_df = damage_df.drop(drop_index)

    #Removes Rows where EV < 60
    damage_df = damage_df.drop(damage_df[damage_df.ExitSpeed < 55.0].index)

    #Switches from Pitcher view to catcher view
    damage_df['PlateLocSide'] = (damage_df['PlateLocSide'] * -1)

    return damage_df

#Creates data Frame for EV Heat map Overhead View
def data_frame_for_overhead_damage_chart():
    global csv_df
    damage_df = csv_df
    damage_df = damage_df.drop(damage_df[damage_df.Batter != names[count]].index)

    #Removes Unnessasry Collumns
    remove_list = damage_df.columns.values.tolist()
    remove_list.remove("Batter")
    #POS Z is essentally -platelocside(already in catchers view)
    remove_list.remove("ContactPositionZ")
    #POS X is distance from front of home plate where contact occured in line with pitcher
    remove_list.remove("ContactPositionX")
    remove_list.remove("BatterSide")
    remove_list.remove("ExitSpeed")
    damage_df = damage_df.drop(remove_list, axis=1)

    #Removes Rows where EV or Contact Position is nan
    drop_index = []
    for index, row in damage_df.iterrows():
        if math.isnan(row['ExitSpeed']) == True:
            drop_index.append(index)
        if math.isnan(row['ContactPositionX']) == True:
            drop_index.append(index)
        if math.isnan(row['ContactPositionZ']) == True:
            drop_index.append(index)
    damage_df = damage_df.drop(drop_index)

    #Removes Rows where EV is under 60
    damage_df = damage_df.drop(damage_df[damage_df.ExitSpeed < 55.0].index)

    return damage_df

#Creates an overhead damage chart with matplotlib and saves the file in sheets
def damage_chart_overhead(damage_df):
    print(damage_df)
    x = []
    y = []
    array = []
    print(damage_df)
    for i in range(17):
        row = []
        for j in range(8):
            temp_df = damage_df
            top_limit = 4.50-0.375*i
            bottom_limit = 4.50-0.375*(i+1)
            left_limit = -1.500 + 0.375*j
            right_limit = -1.500 + 0.375*(j+1)
            if top_limit == 0:
                print(i)
            #drops data from df not in cell needed
            too_left = temp_df[(temp_df['ContactPositionZ'] < left_limit)].index
            temp_df = temp_df.drop(too_left)
            too_right = temp_df[(temp_df['ContactPositionZ'] > right_limit)].index
            temp_df = temp_df.drop(too_right)
            too_high = temp_df[(temp_df['ContactPositionX'] > top_limit)].index
            temp_df = temp_df.drop(too_high)
            too_low = temp_df[(temp_df['ContactPositionX'] < bottom_limit)].index
            temp_df = temp_df.drop(too_low)

            avg_ev_for_zone = round(temp_df["ExitSpeed"].mean(),1)
            
            #if math.isnan(avg_ev_for_zone) == True:
                #avg_ev_for_zone = 60
                

            row.append(avg_ev_for_zone)

        array.append(row)
    df = pd.DataFrame(array)
    print(df)
    df.to_numpy()
    x = np.arange(0, df.shape[1])
    y = np.arange(0, df.shape[0])
    #mask invalid values
    df = np.ma.masked_invalid(df)
    xx, yy = np.meshgrid(x, y)
    #get only the valid values
    x1 = xx[~df.mask]
    y1 = yy[~df.mask]
    newarr = df[~df.mask]

    GD1 = interpolate.griddata((x1, y1), newarr.ravel(), (xx, yy), method='linear', fill_value=55.0)

    df = pd.DataFrame(GD1)

    print(df)

    fig, ax = plt.subplots()
    fig = plt.imshow(df, cmap = 'jet',vmin=55,vmax=95, interpolation='bicubic')
    plt.plot([1.611, 5.389], [11.5, 11.5], color='black', linestyle='-', linewidth=2)
    plt.plot([1.611, 1.611], [11.5, 13.3888889], color='black', linestyle='-', linewidth=2)
    plt.plot([5.389, 5.389], [11.5, 13.3888889], color='black', linestyle='-', linewidth=2)
    plt.plot([5.389, 3.5], [13.3888889, 15.3888889], color='black', linestyle='-', linewidth=2)
    plt.plot([1.61111, 3.61111], [13.3888889, 15.3888889], color='black', linestyle='-', linewidth=2)
    #cbar = plt.colorbar(fig)
    #plt.title("Exit Velo Heat Map")
    plt.axis('off')
    newpath = os.path.join("Sheets", names[count])
    if not os.path.exists(newpath):
        os.makedirs(newpath)
    #saves plot in folder
    plt.savefig(os.path.join("Sheets", names[count], 'overheadheatmap.png'),bbox_inches='tight', pad_inches = 0)

    return

def damage_chart(damage_df):
    x = []
    y = []
    array = []
    print(damage_df)
    for i in range(10):
        row = []
        for j in range(8):
            temp_df = damage_df
            top_limit = 4.267222-0.35819444*i
            bottom_limit = 4.267222-0.35819444*(i+1)
            left_limit = -1.4166667 + 0.35416667*j
            right_limit = -1.416667 + 0.35416667*(j+1)

            #drops data from df not in cell needed
            too_left = temp_df[(temp_df['PlateLocSide'] < left_limit)].index
            temp_df = temp_df.drop(too_left)
            too_right = temp_df[(temp_df['PlateLocSide'] > right_limit)].index
            temp_df = temp_df.drop(too_right)
            too_high = temp_df[(temp_df['PlateLocHeight'] > top_limit)].index
            temp_df = temp_df.drop(too_high)
            too_low = temp_df[(temp_df['PlateLocHeight'] < bottom_limit)].index
            temp_df = temp_df.drop(too_low)

            avg_ev_for_zone = round(temp_df["ExitSpeed"].mean(),1)
            
            #if math.isnan(avg_ev_for_zone) == True:
                #avg_ev_for_zone = 60
                

            row.append(avg_ev_for_zone)

        array.append(row)
    df = pd.DataFrame(array)
    print(df)
    df.to_numpy()
    x = np.arange(0, df.shape[1])
    y = np.arange(0, df.shape[0])
    #mask invalid values
    df = np.ma.masked_invalid(df)
    xx, yy = np.meshgrid(x, y)
    #get only the valid values
    x1 = xx[~df.mask]
    y1 = yy[~df.mask]
    newarr = df[~df.mask]

    GD1 = interpolate.griddata((x1, y1), newarr.ravel(),
                            (xx, yy),
                                method='linear', fill_value=55)

    df = pd.DataFrame(GD1)
    



    #for a in range(6):
    #    if math.isnan(df.iloc[0,a]):
    #        if a > 0 and a < 6:
    #            df.iloc[0,a] = (df.iloc[0,a-1]+df.iloc[0,a+1])




    #df_smooth = gaussian_filter(df, sigma=1)
    #sns.heatmap(df_smooth, cmap='Spectral_r')
    print(df)

    fig, ax = plt.subplots()
    rect = patches.Rectangle((1.5, 1.5), 4, 6, linewidth=1, edgecolor='black', facecolor='none')
    ax.add_patch(rect)
    horizontal_line1 = lines.Line2D([1.5, 5.5], [5.5, 5.5], linewidth=1, color='black')
    horizontal_line2 = lines.Line2D([1.5, 5.5], [3.5, 3.5], linewidth=1, color='black')
    vertical_line1 = lines.Line2D([2.833333, 2.833333], [1.5, 7.5], linewidth=1, color='black')
    vertical_line2 = lines.Line2D([4.166667, 4.166667], [1.5, 7.5], linewidth=1, color='black')
    ax.add_line(horizontal_line1)
    ax.add_line(horizontal_line2)
    ax.add_line(vertical_line1)
    ax.add_line(vertical_line2)
    fig = plt.imshow(df, cmap = 'jet',vmin=55,vmax=95,  interpolation='bicubic')
    #cbar = plt.colorbar(fig)
    fig.axes.get_xaxis().set_visible(False)
    fig.axes.get_yaxis().set_visible(False)
    #plt.title("Exit Velo Heat Map")
    newpath = os.path.join("Sheets", names[count])
    if not os.path.exists(newpath):
        os.makedirs(newpath)
    #saves plot in folder
    plt.axis('off')
    plt.savefig(os.path.join("Sheets", names[count], 'heatmap.png'), bbox_inches='tight', pad_inches = 0)



    

    
    #print(df_smooth)









    return

def presentation (tabledata):
    prs = Presentation("KnightsTemplate.pptx")
    prs.slide_width = Inches(11)
    prs.slide_height = Inches(8.5)
    #slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide = prs.slides.get(257)

    name = names[count].split()
    print(name)
    full_name = name[-1] + ' ' + name[0]
    full_name = full_name[:-1]
    title = slide.shapes.title
    title.text = full_name
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.name = 'Avenir Next LT Pro Light'
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ---add table to slide---
    x, y, cx, cy = Inches(0.15), Inches(0.95), Inches(10.7), Inches(1)
    shape = slide.shapes.add_table(2, 9, x, y, cx, cy)
    table = shape.table

    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{073A0DAA-6AF3-43AB-8588-CEC1D06C72B9}'
    tbl[0][-1].text = style_id

    #creating labels for values in all tables
    cell = table.cell(0, 0)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text = 'AVG EV'
    cell.text_frame.paragraphs[0].font.name = 'Bahnschrift Condensed'
    cell.text_frame.paragraphs[0].font.size = Pt(22)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    cell = table.cell(0,1)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text = 'MAX EV'
    cell.text_frame.paragraphs[0].font.name = 'Bahnschrift Condensed'
    cell.text_frame.paragraphs[0].font.size = Pt(22)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    cell = table.cell(0,2)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text = "Swing %"
    cell.text_frame.paragraphs[0].font.name = 'Bahnschrift Condensed'
    cell.text_frame.paragraphs[0].font.size = Pt(22)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER    

    cell = table.cell(0,3)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text = "Chase %"
    cell.text_frame.paragraphs[0].font.name = 'Bahnschrift Condensed'
    cell.text_frame.paragraphs[0].font.size = Pt(22)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  

    cell = table.cell(0,4)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text = "InZoneSwing %"
    cell.text_frame.paragraphs[0].font.name = 'Bahnschrift Condensed'
    cell.text_frame.paragraphs[0].font.size = Pt(14)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  

    cell = table.cell(0,5)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text = "Strikeout %"
    cell.text_frame.paragraphs[0].font.name = 'Bahnschrift Condensed'
    cell.text_frame.paragraphs[0].font.size = Pt(18)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  

    cell = table.cell(0,6)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text = "BB+HBP/K"
    cell.text_frame.paragraphs[0].font.name = 'Bahnschrift Condensed'
    cell.text_frame.paragraphs[0].font.size = Pt(22)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  

    cell = table.cell(0,7)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text = "BABIP"
    cell.text_frame.paragraphs[0].font.name = 'Bahnschrift Condensed'
    cell.text_frame.paragraphs[0].font.size = Pt(22)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    cell = table.cell(0,8)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text = "wOBA"
    cell.text_frame.paragraphs[0].font.name = 'Bahnschrift Condensed'
    cell.text_frame.paragraphs[0].font.size = Pt(22)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


    #nested for loop which puts vales from averages into table
    for j in range(9):
        cell = table.cell((1),j)
        cell.text = tabledata[j]
        cell.text_frame.paragraphs[0].font.name = 'Bahnschrift Condensed'
        cell.text_frame.paragraphs[0].font.size = Pt(30)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE 
        
    prs.save(os.path.join("Sheets", names[count], names[count] + '.pptx'))

    x, y, cx, cy = Inches(0.15), Inches(2.15), Inches(10.7), Inches(1)
    shape2 = slide.shapes.add_table(2, 15, x, y, cx, cy)
    table2 = shape2.table

    tbl =  shape2._element.graphic.graphicData.tbl
    style_id = '{073A0DAA-6AF3-43AB-8588-CEC1D06C72B9}'
    tbl[0][-1].text = style_id
    text = ['PA', 'AB', 'H', '2B', '3B', 'HR', 'RBI','R', 'BB', 'HBP', 'K', 'AVG', 'OBP','SLG', 'OPS']
    #creating labels for values in all tables
    for i in range(15):
        cell = table2.cell(0, i)
        cell.text = text[i]
        cell.text_frame.paragraphs[0].font.size = Pt(22)
        cell.text_frame.paragraphs[0].font.name = 'Bahnschrift Condensed'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    for j in range(15):
        cell = table2.cell((1),j)
        cell.text = tabledata[j+9]
        cell.text_frame.paragraphs[0].font.name = 'Bahnschrift Condensed'
        if j > 9:
            cell.text_frame.paragraphs[0].font.size = Pt(22)
        else:
            cell.text_frame.paragraphs[0].font.size = Pt(24)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE 

    img_swing_chart = Image.open(os.path.join("Sheets", names[count], 'swingchart.png'))


    # Get the size of the original image
    width, height = img_swing_chart.size

    # Set the left and right margins to crop
    left = 23
    right = width - 23

    # Crop the image
    cropped_image = img_swing_chart.crop((left, 0, right, height))

    # Save the cropped image
    cropped_image.save(os.path.join("Sheets", names[count], 'swingchart.png'))
    img_swing_chart = os.path.join("Sheets", names[count], 'swingchart.png')
    shape = slide.shapes.add_picture(img_swing_chart,Inches(0.15),Inches(3.77), width=Inches(2.91), height=Inches(3.44))
    line = shape.line
    line.color.rgb = RGBColor(117, 117, 117)
    line.width = Inches(0.05)
    
    prs.save(os.path.join("Sheets", names[count], names[count] + '.pptx'))

    img_heat_map = os.path.join("Sheets", names[count], 'heatmap.png')
    shape = slide.shapes.add_picture(img_heat_map,Inches(3.22),Inches(3.77), width=Inches(2.75), height=Inches(3.44))
    line = shape.line
    line.color.rgb = RGBColor(117, 117, 117)
    line.width = Inches(0.05)

    prs.save(os.path.join("Sheets", names[count], names[count] + '.pptx'))

    img_heat_map_overhead = os.path.join(os.path.join("Sheets", names[count], 'overheadheatmap.png'))
    shape = slide.shapes.add_picture(img_heat_map_overhead,Inches(6.12),Inches(3.77), width=Inches(1.61), height=Inches(3.44))
    line = shape.line
    line.color.rgb = RGBColor(117, 117, 117)
    line.width = Inches(0.05)

    prs.save(os.path.join("Sheets", names[count], names[count] + '.pptx'))

    scatter = os.path.join(os.path.join("Sheets", names[count], 'zone' + '.png'))
    shape = slide.shapes.add_picture(scatter,Inches(7.88),Inches(3.77), width=Inches(2.97), height=Inches(3.44))
    line = shape.line
    line.color.rgb = RGBColor(117, 117, 117)
    line.width = Inches(0.05)

    prs.save(os.path.join("Sheets", names[count], names[count] + '.pptx'))

def pitch_strike_called_df():
    global csv_df
    player_df = csv_df

    #Drops all rows where player isn't hitting
    #Drops all rows where player didn't swing
    player_df = player_df.drop(player_df[player_df.PitchCall == 'BallCalled'].index)
    player_df = player_df.drop(player_df[player_df.PitchCall == 'StrikeSwinging'].index)
    player_df = player_df.drop(player_df[player_df.PitchCall == 'FoulBall'].index)
    player_df = player_df.drop(player_df[player_df.PitchCall == 'BallinDirt'].index)
    player_df = player_df.drop(player_df[player_df.PitchCall == 'HitByPitch'].index)
    player_df = player_df.drop(player_df[player_df.PitchCall == 'InPlay'].index)
    #Creating a list of collums to remove from DF and then removing collums we need from list
    remove_list = player_df.columns.values.tolist()
    remove_list.remove("Batter")
    remove_list.remove("PlateLocHeight")
    remove_list.remove("PlateLocSide")
    remove_list.remove("PitchCall")
    remove_list.remove("BatterSide")
    #Removes all collums other than those with .remove above from data frame
    player_df = player_df.drop(remove_list, axis=1)
    #Switches plate loc side values to be in catcher view by multiplying by -1
    player_df['PlateLocSide'] = (player_df['PlateLocSide'] * -1)
    print(player_df.to_string())


    return player_df

def pitch_loc_chart(player_df):
    #Pulls image for background will have to imput if statemnt to deptermine right vs left
    rhhs = ['Burke, Isaiah','Cedillo, Ruben']
    img = plt.imread("RHH.png")
    fig, ax = plt.subplots(figsize=(6, 6))
    sns.set_style("white")
    #Creates density plot, camp is color scheme and alpha is transperacy
    chart = plt.scatter(x=player_df.PlateLocSide, y=player_df.PlateLocHeight, )
    #creates demenstions for graph plus displays image
    ax.imshow(img, extent=[-2.63,2.665,-0.35,5.30], aspect=1)
      # remove the ticks
    #rect = patches.Rectangle((-0.708333, 1.6466667), 1.4166667, 1.90416667, linewidth=1, edgecolor='black', facecolor='none')
    #ax.add_patch(rect)
    #creates path for plot to be saved in there isn't one
    newpath = os.path.join("Sheets", names[count])
    if not os.path.exists(newpath):
        os.makedirs(newpath)
    #saves plot in folder
    #plt.savefig(os.path.join("Sheets", names[0], 'swingchart.png'))
    plt.show()

    return

def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()


def pitchResult(player_df):
    plt.clf()
    player_df['PlateLocSide'] = player_df['PlateLocSide']*-1
    strikeC,strikeS,ball,foul,inPlay = 0,0,0,0,0
    fig, ax = plt.subplots(figsize=(4.2878, 5))
    ax.set_xlim(-2.1439, 2.1439)
    ax.set_ylim(0, 5)
    for j in range(len(player_df.index)):
    #if the player is on OSU and it is the batter the function is asking for
        if ((player_df.loc[j, 'PitchCall']=='StrikeCalled')):
            ax.plot(player_df['PlateLocSide'][j], player_df['PlateLocHeight'][j], markersize=5,color='Black',markerfacecolor='#ff726f', marker='o', linestyle='')
            strikeC+=1
        #stike swinging
        elif (player_df.loc[j, 'PitchCall']=='StrikeSwinging'):
            ax.plot(player_df['PlateLocSide'][j], player_df['PlateLocHeight'][j], markersize=5,color='Black',markerfacecolor='#8B0000', marker='o', linestyle='')
            strikeS+=1
        #in play
        elif (player_df.loc[j, 'PitchCall']=='InPlay'):
            ax.plot(player_df['PlateLocSide'][j], player_df['PlateLocHeight'][j], markersize=5,color='Black',markerfacecolor='Blue', marker='o', linestyle='')
            inPlay+=1
        #foul ball
        elif (player_df.loc[j, 'PitchCall']=='FoulBall'):
            ax.plot(player_df['PlateLocSide'][j], player_df['PlateLocHeight'][j], markersize=5,color='Black',markerfacecolor='#FFB52E', marker='o', linestyle='')
            foul+=1
        #ball
        elif (player_df.loc[j, 'PitchCall'] == 'BallCalled'):
            ax.plot(player_df['PlateLocSide'][j], player_df['PlateLocHeight'][j], markersize=5,color='Black',markerfacecolor='Green', marker='o', linestyle='')
            ball+=1

    ax.set_xlim(-2.1439, 2.1439)
    ax.set_ylim(0, 5)

    # Add legend
    legend_elements = [
    plt.Line2D([0], [0], marker='o', color='w', label='Strike Called: ' + str(strikeC), markerfacecolor='#ff726f', markersize=5),
    plt.Line2D([0], [0], marker='o', color='w', label='Strike Swinging: ' + str(strikeS), markerfacecolor='#8B0000', markersize=5),
    plt.Line2D([0], [0], marker='o', color='w', label='In Play: ' + str(inPlay), markerfacecolor='Blue', markersize=5),
    plt.Line2D([0], [0], marker='o', color='w', label='Foul Ball: ' + str(foul), markerfacecolor='#FFB52E', markersize=5),
    plt.Line2D([0], [0], marker='o', color='w', label='Ball Called: '+ str(ball), markerfacecolor='Green', markersize=5)
]

    rect = plt.Rectangle((-0.708333, 1.6466667), 1.4166667, 1.90416667, linewidth=1, edgecolor='black', facecolor='none', zorder=2)
    plt.gca().add_patch(rect)


    print("Strike swing: "+str(strikeS))
    print("Strike called: "+str(strikeC))
    print("Foul: "+str(foul))
    print("In play: "+ str(inPlay))
    print("Ball: "+str(ball))

    plt.axis('off')
    ax.legend(handles=legend_elements, loc='upper left', fontsize='small')
    plt.savefig(os.path.join("Sheets", names[count], 'zone' + '.png'), bbox_inches='tight', pad_inches = 0)

def get_stats():
    baseball_cube = pd.read_html("http://pointstreak.com/baseball/team_stats.html?teamid=160840&seasonid=33635")
    hitters = baseball_cube[1]
    # Set the first row as the new column headers
    #hitters.columns = hitters.iloc[0]

    # Remove the first row from the DataFrame
    #hitters = hitters[1:].reset_index(drop=True)

    return hitters

def get_player_stats(name, df):
    print(df.columns)
    # Filter the DataFrame based on the 'player name' column
    filtered_df = df[df['Player'] == name]
    filtered_df = filtered_df.reset_index(drop=True)

    print(filtered_df)

    return filtered_df

'''def main():
    global count
    global names
    global csv_df
    #names = csv_df.loc[csv_df['BatterTeam'] == 'ORE_BEA', 'Batter'].unique()
    #names = names.tolist()
    for i in range(len(names)):
        name = names[count].split()
        full_name = name[-1] + ' ' + name[0]
        full_name = full_name[:-1]
        player_df = get_player_df()
        #all_results(player_df)
        stats = get_player_stats(full_name,get_stats())
        swing2d_density_plot(csv_to_swing_df())
        damage_chart(data_frame_for_damage_chart())
        damage_chart_overhead(data_frame_for_overhead_damage_chart())
        pitchResult(player_df)
        presentation(find_table_metrics_using_online(stats))
    
        
        #PPTtoPDF(os.path.join("C:\\Users\\cmhea\\OneDrive\\Documents\\baseball\\HittersReports\\HittersReport\\Sheets", names[count], names[count] + '.pptx'),os.path.join("C:\\Users\\cmhea\\OneDrive\\Documents\\baseball\\HittersReports\\HittersReport\\Sheets\\PDF", names[count] + '.pdf'))
        count = count + 1'''
#pitch_loc_chart(pitch_strike_called_df())

def main():
    global count
    global names
    global csv_df
    #names = csv_df.loc[csv_df['BatterTeam'] == 'ORE_BEA', 'Batter'].unique()
    #names = names.tolist()
    for i in range(len(names)):
        # Splitting the string into first name and last name
        last_name, first_name = names[i].split(', ')
        # Extracting the first letter of the first name
        first_initial = first_name[0]
        # Combining the first initial with the last name
        new_name = f"{last_name}, {first_initial}"
        player_df = get_player_df()
        #all_results(player_df)
        stats = get_player_stats(new_name,get_stats())
        swing2d_density_plot(csv_to_swing_df())
        damage_chart(data_frame_for_damage_chart())
        damage_chart_overhead(data_frame_for_overhead_damage_chart())
        pitchResult(player_df)
        presentation(find_table_metrics_using_online(stats))
    
        
        #PPTtoPDF(os.path.join("C:\\Users\\cmhea\\OneDrive\\Documents\\baseball\\HittersReports\\HittersReport\\Sheets", names[count], names[count] + '.pptx'),os.path.join("C:\\Users\\cmhea\\OneDrive\\Documents\\baseball\\HittersReports\\HittersReport\\Sheets\\PDF", names[count] + '.pdf'))
        count = count + 1

main()